from __future__ import annotations

import csv
import re
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK, WD_TAB_ALIGNMENT, WD_TAB_LEADER
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


REPO_ROOT = Path(__file__).resolve().parents[1]
PROJECT_ROOT = REPO_ROOT.parent
REPORT_OUTPUTS = PROJECT_ROOT / "report_outputs"
REPORT_ASSETS = PROJECT_ROOT / "report_assets"
COLAB_RESULTS = PROJECT_ROOT / "colab_results"
DELIVERABLES = PROJECT_ROOT / "final_deliverables"
REPORT_TEMPLATE = PROJECT_ROOT / "myfriendsproject" / "04-4902-Project_Report_Template.docx"
GROUPED_SUMMARY = PROJECT_ROOT / "upload_to_drive" / "underwater_resnet_project" / "experiments" / "grouped_v1" / "summaries" / "metric_summary_mean_std.csv"
def main() -> None:
    if not GROUPED_SUMMARY.exists():
        raise FileNotFoundError(
            "Grouped 3-model x 3-seed results are not available. The existing report contains "
            "superseded random-split results; complete scripts/run_grouped_experiments.py before regenerating deliverables. "
            f"Expected: {GROUPED_SUMMARY}"
        )
    raise RuntimeError(
        "Grouped results exist, but this generator still contains superseded random-split narrative. "
        "Update the report tables and conclusions from metric_summary_mean_std.csv before generating submission files."
    )
    DELIVERABLES.mkdir(parents=True, exist_ok=True)
    REPORT_OUTPUTS.mkdir(parents=True, exist_ok=True)
    create_compact_comparison_assets()
    create_docx_report()
    create_presentation_outline()
    create_presentation_pptx()
    create_submission_readme()
    print(f"Final deliverables written to: {DELIVERABLES}")


def create_docx_report() -> None:
    doc = Document(REPORT_TEMPLATE) if REPORT_TEMPLATE.exists() else Document()
    replace_template_cover(doc)
    rebuild_template_contents(doc)
    remove_template_body(doc)
    setup_document(doc)
    add_summary_section(doc)
    add_summary(doc)
    add_body_section(doc)
    add_introduction(doc)
    add_background(doc)
    add_system_requirements(doc)
    add_system_architecture(doc)
    add_dataset_and_development(doc)
    add_results_and_evaluation(doc)
    add_conclusions(doc)
    add_references(doc)

    output = DELIVERABLES / "final_report_150210321.docx"
    try:
        doc.save(output)
        print(f"Report written to: {output}")
    except PermissionError:
        fallback = DELIVERABLES / "final_report_150210321_template.docx"
        doc.save(fallback)
        print(f"Report was locked; fallback written to: {fallback}")


def replace_template_cover(doc: Document) -> None:
    replacements = {
        "Building A Large Language Model Based Artificial General Intelligence (AGI) System":
            "Enhancing Diffusion-Based Underwater Image Restoration with Residual Networks",
        "15010000 Elon MUSK": "150210321 Bedirhan OZTURK",
        "15010001 Bill GATES": "",
        "Assoc. Prof. Dr. A. Cüneyd TANTUĞ": "Prof. Dr. Behcet Ugur Toreyin",
        "May, 2024": "August, 2026",
    }
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text in replacements:
            set_paragraph_text(paragraph, replacements[text], keep_style=True)


def rebuild_template_contents(doc: Document) -> None:
    contents_index = next(
        (i for i, p in enumerate(doc.paragraphs) if p.text.strip() == "CONTENTS"), None
    )
    summary_index = next(
        (i for i, p in enumerate(doc.paragraphs) if p.text.strip() == "SUMMARY"), None
    )
    if contents_index is None or summary_index is None:
        return

    contents_heading = doc.paragraphs[contents_index]
    set_paragraph_text(contents_heading, "CONTENTS")
    contents_heading.style = doc.styles["Normal"]
    contents_heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
    contents_heading.paragraph_format.left_indent = Inches(0.5)
    contents_heading.paragraph_format.space_after = Pt(6)
    contents_run = contents_heading.runs[0]
    contents_run.font.name = "Calibri"
    contents_run.font.size = Pt(18)
    contents_run.bold = False
    for paragraph in list(doc.paragraphs[contents_index + 1:summary_index]):
        # Keep section-break paragraphs because they preserve the official cover/contents/body layout.
        if paragraph._p.pPr is not None and paragraph._p.pPr.sectPr is not None:
            continue
        paragraph._element.getparent().remove(paragraph._element)

    entries = [
        ("1", "INTRODUCTION", "1"),
        ("2", "BACKGROUND", "4"),
        ("2.1", "Underwater Image Formation and Degradation", "4"),
        ("2.2", "Underwater Restoration Methods", "4"),
        ("2.3", "U-Net Backbones", "5"),
        ("2.4", "Residual Learning", "6"),
        ("2.5", "Diffusion Models for Restoration", "6"),
        ("2.6", "Image Quality Metrics", "7"),
        ("3", "SYSTEM REQUIREMENTS", "9"),
        ("3.1", "Design Constraints and Relevant Engineering Standards", "9"),
        ("3.2", "Functional Requirements", "9"),
        ("3.3", "Non-Functional Requirements", "10"),
        ("3.4", "Evaluation Methodology", "11"),
        ("4", "SYSTEM ARCHITECTURE", "12"),
        ("4.1", "Dataset and Pairing Pipeline", "12"),
        ("4.2", "Conditional Diffusion Pipeline", "13"),
        ("4.3", "Baseline Conditional U-Net", "14"),
        ("4.4", "Proposed Residual Backbone", "14"),
        ("5", "DATASET PREPARATION AND MODEL DEVELOPMENT", "16"),
        ("5.1", "Dataset Summary", "16"),
        ("5.2", "Preprocessing and Fixed Splits", "17"),
        ("5.3", "Training Configuration", "18"),
        ("5.4", "Model Development Timeline", "18"),
        ("6", "RESULTS AND EVALUATION", "20"),
        ("6.1", "Main Quantitative Results", "20"),
        ("6.2", "Default Baseline vs Residual", "21"),
        ("6.3", "Parameter-Matched Capacity Control", "23"),
        ("6.4", "Training Duration Ablation", "24"),
        ("6.5", "Resolution Ablation", "25"),
        ("6.6", "Qualitative and Cross-Dataset Scope", "27"),
        ("6.7", "Limitations", "28"),
        ("7", "CONCLUSIONS AND FUTURE WORKS", "30"),
        ("8", "REFERENCES", "32"),
    ]
    anchor = contents_heading._element
    for number, title, page in entries:
        paragraph = OxmlElement("w:p")
        ppr = OxmlElement("w:pPr")
        style = OxmlElement("w:pStyle")
        style.set(qn("w:val"), "TOC2" if "." in number else "TOC1")
        ppr.append(style)
        spacing = OxmlElement("w:spacing")
        spacing.set(qn("w:before"), "0")
        spacing.set(qn("w:after"), "0")
        if "." not in number:
            spacing.set(qn("w:before"), "120")
        ppr.append(spacing)
        indent = OxmlElement("w:ind")
        if "." in number:
            indent.set(qn("w:left"), "1680")
            indent.set(qn("w:hanging"), "720")
        else:
            indent.set(qn("w:left"), "1200")
            indent.set(qn("w:hanging"), "480")
        ppr.append(indent)
        tabs = OxmlElement("w:tabs")
        tab = OxmlElement("w:tab")
        tab.set(qn("w:val"), "right")
        tab.set(qn("w:leader"), "dot")
        tab.set(qn("w:pos"), "10080")
        tabs.append(tab)
        ppr.append(tabs)
        paragraph.append(ppr)
        run = OxmlElement("w:r")
        rpr = OxmlElement("w:rPr")
        font = OxmlElement("w:rFonts")
        font.set(qn("w:ascii"), "Calibri")
        font.set(qn("w:hAnsi"), "Calibri")
        rpr.append(font)
        size = OxmlElement("w:sz")
        size.set(qn("w:val"), "24")
        rpr.append(size)
        if "." not in number:
            rpr.append(OxmlElement("w:b"))
        run.append(rpr)
        text = OxmlElement("w:t")
        text.text = f"{number}  {title.upper()}"
        run.append(text)
        tab_run = OxmlElement("w:r")
        tab_element = OxmlElement("w:tab")
        tab_run.append(tab_element)
        page_run = OxmlElement("w:r")
        page_rpr = OxmlElement("w:rPr")
        page_font = OxmlElement("w:rFonts")
        page_font.set(qn("w:ascii"), "Calibri")
        page_font.set(qn("w:hAnsi"), "Calibri")
        page_rpr.append(page_font)
        page_size = OxmlElement("w:sz")
        page_size.set(qn("w:val"), "24")
        page_rpr.append(page_size)
        page_run.append(page_rpr)
        page_text = OxmlElement("w:t")
        page_text.text = page
        page_run.append(page_text)
        paragraph.extend([run, tab_run, page_run])
        anchor.addnext(paragraph)
        anchor = paragraph


def remove_template_body(doc: Document) -> None:
    start = next((i for i, p in enumerate(doc.paragraphs) if p.text.strip() == "SUMMARY"), None)
    if start is None:
        doc.add_page_break()
        return
    for paragraph in list(doc.paragraphs[start:]):
        paragraph._element.getparent().remove(paragraph._element)


def setup_document(doc: Document) -> None:
    styles = doc.styles
    styles["Normal"].font.name = "Calibri"
    styles["Normal"].font.size = Pt(12)
    if "Body Text" in styles:
        styles["Body Text"].font.name = "Calibri"
        styles["Body Text"].font.size = Pt(12)
    styles["Heading 1"].font.name = "Arial"
    styles["Heading 1"].font.size = Pt(16)
    styles["Heading 1"].font.bold = True
    styles["Heading 1"].font.color.rgb = RGBColor(0, 0, 0)
    styles["Heading 2"].font.name = "Arial"
    styles["Heading 2"].font.size = Pt(14)
    styles["Heading 2"].font.bold = True
    styles["Heading 2"].font.color.rgb = RGBColor(0, 0, 0)
    title_style = styles["Title"] if "Title" in styles else styles.add_style("Title", WD_STYLE_TYPE.PARAGRAPH)
    title_style.font.name = "Arial"
    title_style.font.size = Pt(22)
    title_style.font.bold = True

    # Section 0 remains the official cover, including the official blue header and logo footer.
    for index, section in enumerate(doc.sections):
        if index == 0:
            continue
        configure_report_section(section, numbered=True, preamble=True)
        set_page_numbering(section, "lowerRoman", 1)


def add_body_section(doc: Document) -> None:
    section = doc.add_section(WD_SECTION.NEW_PAGE)
    configure_report_section(section, numbered=True, preamble=False)
    set_page_numbering(section, "decimal", 1)


def add_summary_section(doc: Document) -> None:
    section = doc.add_section(WD_SECTION.NEW_PAGE)
    configure_report_section(section, numbered=True, preamble=True)
    set_page_numbering(section, "lowerRoman", 2)


def configure_report_section(section, numbered: bool, preamble: bool) -> None:
    section.left_margin = Inches(0.5)
    section.right_margin = Inches(0.5)
    section.top_margin = Inches(1.26 if preamble else 1.03)
    section.bottom_margin = Inches(0.65 if preamble else 0.36)
    section.header_distance = Inches(0)
    section.footer_distance = Inches(0.52 if preamble else 0.23)
    section.header.is_linked_to_previous = False
    section.footer.is_linked_to_previous = False
    clear_container(section.header)
    clear_container(section.footer)
    if numbered:
        footer_p = section.footer.paragraphs[0]
        footer_p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        add_page_field(footer_p)


def set_page_numbering(section, number_format: str, start: int) -> None:
    for existing in section._sectPr.xpath("./w:pgNumType"):
        existing.getparent().remove(existing)
    page_numbering = OxmlElement("w:pgNumType")
    page_numbering.set(qn("w:fmt"), number_format)
    page_numbering.set(qn("w:start"), str(start))
    section._sectPr.append(page_numbering)


def clear_container(container) -> None:
    for paragraph in container.paragraphs:
        paragraph.clear()
    for table in list(container.tables):
        table._element.getparent().remove(table._element)
    for drawing in list(container._element.xpath(".//w:drawing")):
        drawing.getparent().remove(drawing)


def add_page_field(paragraph) -> None:
    run = paragraph.add_run()
    run.font.name = "Calibri"
    run.font.size = Pt(9)
    begin = OxmlElement("w:fldChar")
    begin.set(qn("w:fldCharType"), "begin")
    instruction = OxmlElement("w:instrText")
    instruction.set(qn("xml:space"), "preserve")
    instruction.text = " PAGE "
    end = OxmlElement("w:fldChar")
    end.set(qn("w:fldCharType"), "end")
    run._r.extend([begin, instruction, end])


def add_summary(doc: Document) -> None:
    preamble_heading(doc, "SUMMARY")
    paragraphs(doc, [
        "Underwater imaging supports inspection, observation, robotics, scientific documentation, and monitoring, but the water medium changes the signal before it reaches a camera. Wavelength-dependent absorption removes color components at different rates, suspended particles scatter light, and turbidity lowers contrast while obscuring boundaries. This project addresses a supervised version of that restoration problem: given a turbid underwater color-patch image, produce an image that approaches its paired clear reference. The work does not claim to invert every physical process in arbitrary ocean scenes. Instead, it evaluates whether a residual denoising architecture is useful inside one controlled conditional diffusion pipeline.",
        "The dataset contains 3,672 filename-matched turbid and clear pairs. A fixed split allocates 2,937 pairs to training, 367 to validation, and 368 to testing. Images are converted to RGB, resized to the configured square resolution, converted to tensors, and normalized to [-1, 1]. During diffusion training, Gaussian noise is added to the clear target at a randomly selected one of 100 timesteps. The denoiser receives the noisy target, the turbid condition, and a sinusoidal timestep representation, then predicts the sampled noise. Mean squared noise-prediction error is optimized. This formulation and the data protocol remain constant when the denoising backbone changes.",
        "The default baseline is a two-level conditional U-Net with base width 32. The proposed model uses the same encoder-decoder scale structure but replaces ordinary convolution blocks with residual blocks and includes a second residual bottleneck block. The resulting models have approximately 527,363 and 886,371 trainable parameters, respectively. Because this difference could confound an architectural comparison, a base-width-42 U-Net with approximately 901,797 parameters is evaluated as a capacity control. Additional experiments examine 100 rather than 50 epochs and 256 rather than 128 pixel inputs. Every reported metric is measured on the same 368 paired test items.",
        "At 128 pixels and 50 epochs, the residual model improves the default U-Net from 0.039444 to 0.035727 MSE, from 15.681501 to 16.304393 dB PSNR, from 0.610758 to 0.788963 SSIM, and from 29.645203 to 26.840257 CIE76 Delta E. The parameter-matched U-Net, however, records slightly better MSE, MAE, and PSNR than the residual model, while the residual model remains better in SSIM, Delta E, and output entropy. The defensible conclusion is therefore metric-specific: residual processing improves structural similarity and color-difference behavior over both the default baseline and the capacity control in this experiment, but increased U-Net capacity is competitive or better for direct pixel error. External scenes are shown only as input-condition examples; no cross-dataset restoration claim is made.",
    ])


def add_introduction(doc: Document) -> None:
    heading(doc, "1 INTRODUCTION")
    paragraphs(doc, [
        "A camera used in air observes light after it has interacted with surfaces and traveled through a comparatively transparent medium. Underwater, the propagation path itself becomes a major part of image formation. Water absorbs long wavelengths strongly, suspended particles redirect light, and illumination varies with depth, distance, and local geometry. Consequently, an underwater photograph can exhibit a blue-green cast, muted warm colors, low global contrast, reduced local edge contrast, haze-like veiling, and spatially varying degradation at the same time. These effects are not merely cosmetic. They can make boundaries, markings, instruments, organisms, and man-made structures harder to inspect visually and harder to process with later computer-vision systems.",
        "Underwater image enhancement and restoration are related but distinct objectives. Enhancement seeks a visually preferable result and can operate without a known clean target. Restoration is more naturally framed as estimating a less degraded observation from a measurement and, where available, a reference. The present project uses paired supervision: each turbid color-patch photograph is associated by filename with a corresponding clearer photograph. This design gives a measurable target and permits full-reference metrics. It also narrows the scope. The learned relationship is grounded in the paired acquisition represented by this dataset and should not be described as a universal physical solution for all water types, cameras, depths, or natural scenes.",
        "The engineering question is whether the denoising backbone of a conditional diffusion model benefits from explicit residual learning. Denoising diffusion probabilistic models learn a reverse process by predicting noise applied through a known forward schedule [1]. A conditional restoration variant can make that prediction while observing a degraded image. U-Net architectures are a natural baseline because their contracting path, expanding path, and skip connections combine broad context with spatial detail [2]. Residual networks provide another form of shortcut: a block learns a correction relative to an identity or projected signal, which can ease optimization and preserve information [3]. This project places those ideas in the same implementation and evaluates their consequences under a shared protocol.",
        "The core comparison uses 128 by 128 inputs, 50 epochs, a batch size of 16, a learning rate of 0.0001, seed 42, and a 100-step linear diffusion schedule. The default U-Net and residual model both receive six channels formed by concatenating a three-channel noisy target with the three-channel turbid condition. Their output is a three-channel noise estimate. Test results cover MSE, MAE, PSNR, SSIM, CIE76 Delta E, and entropy. The report gives special attention to the parameter-matched U-Net because the proposed model contains more parameters than the default baseline. Without this control, an apparent residual benefit might simply reflect additional capacity.",
    ])
    add_picture(doc, REPORT_ASSETS / "summary_options" / "option_7_three_condition_color_chart.png", 6.7)
    caption(doc, "Figure 1-1: Real color-chart observations under three acquisition conditions; these are source-condition images, not model outputs")
    paragraphs(doc, [
        "Figure 1-1 establishes the visual context using actual condition photographs. Differences among no-water, clear-water, and turbid-water observations expose loss of saturation, shifts in patch appearance, and reduced separability between neighboring colors. The figure must not be mistaken for a restoration sequence: no baseline or proposed-model output is present. It is included early because a metric-only account can hide the physical character of the input. The controlled chart also explains why a color-space difference measure complements pixel loss and structural similarity in the evaluation.",
    ])
    add_picture(doc, REPORT_ASSETS / "summary_options" / "option_8_three_condition_usaf.png", 6.7)
    caption(doc, "Figure 1-2: Real USAF resolution-target observations under three acquisition conditions; these are source-condition images, not model outputs")
    paragraphs(doc, [
        "The USAF target in Figure 1-2 illustrates the structural side of degradation. Bars that are distinct under a favorable condition become less separable as scattering and turbidity suppress local contrast. Again, these panels document acquisition conditions rather than generated restorations. Their role is to connect the numerical use of SSIM with a practical concern: a useful restoration should retain or recover meaningful spatial organization rather than merely shift average color or reduce a pixel-wise loss.",
        "The project contribution is a controlled implementation study rather than a claim of state-of-the-art performance. It provides a reproducible paired-data pipeline, a conditional diffusion baseline, a residual alternative, a parameter-count control, duration and resolution ablations, quantitative evaluation on a fixed test set, and cautious qualitative examination outside the paired distribution. The main finding is that the residual backbone is clearly preferable to the small default baseline on the selected structural and color-oriented measurements. The capacity control qualifies that statement by showing that a wider U-Net can lead on MSE, MAE, and PSNR. This distinction between architectural behavior and raw capacity is central to the report.",
        "The remainder of the report follows eight sections. Section 2 reviews underwater formation, restoration families, U-Net and residual architectures, diffusion restoration, and the selected metrics. Section 3 defines requirements, constraints, and evaluation rules. Section 4 describes the system architecture. Section 5 records dataset preparation, preprocessing, training, and model-development stages. Section 6 presents all measured results, ablations, qualitative evidence, and limitations. Section 7 states conclusions and realistic future work, while Section 8 lists the literature used to support the technical discussion.",
    ])


def add_background(doc: Document) -> None:
    heading(doc, "2 BACKGROUND")
    subheading(doc, "2.1 Underwater Image Formation and Degradation")
    paragraphs(doc, [
        "A simplified underwater image-formation account separates direct transmission from path radiance. Light reflected by an object is attenuated along the camera-object path, while light scattered toward the camera adds a veiling component. The transmission depends on distance and on an attenuation coefficient that varies by wavelength. As distance or turbidity increases, the direct component weakens and the scattered component can dominate. This resembles atmospheric haze models at a high level, but underwater attenuation is notably wavelength dependent and strongly affected by the composition of the water. Red light generally disappears sooner than blue or green light, causing the familiar cool cast and making color recovery ill posed when little red-channel information remains [5], [6].",
        "Absorption and scattering produce different visible consequences. Absorption removes energy and changes the balance among channels. Forward scattering spreads light near its original path and softens detail, while backscatter contributes light unrelated to the desired surface signal and lowers contrast. Non-uniform artificial lighting can add bright central regions, falloff, and shadows. Camera response, white balance, sensor noise, and compression are additional factors. A single observed RGB value may therefore be compatible with several combinations of surface reflectance, distance, illumination, water type, and camera processing. Restoration cannot uniquely determine all of those latent quantities from one image without assumptions or learned prior information.",
        "The paired chart data used here captures a repeatable relationship between turbid and clearer observations. It is valuable because corresponding spatial regions and known color patches make full-reference evaluation possible. At the same time, chart patches differ from broad natural scenes containing vegetation, animals, sediment, specular surfaces, and depth variation. A network can learn dataset-specific statistics such as the colors, texture scale, and acquisition setup. This is why the report treats the fixed test split as evidence for in-distribution paired restoration and does not automatically extend the same numerical conclusion to unrelated underwater photographs.",
        "The term clear reference also requires care. A paired reference is the target available in the dataset, not a proof of the object's exact radiance in an ideal physical environment. Differences in acquisition conditions may include illumination and camera effects as well as water clarity. The supervised task is therefore operational: reproduce the paired clear-domain appearance sufficiently well according to complementary metrics. This operational definition supports controlled comparison of two denoisers even though it does not identify a complete physical inverse model.",
    ])
    subheading(doc, "2.2 Underwater Restoration Methods")
    paragraphs(doc, [
        "Underwater methods can be grouped broadly into non-learning enhancement, physics-informed restoration, and learned image-to-image approaches. Histogram manipulation, contrast stretching, white balancing, gamma correction, and fusion can improve visibility without explicit paired training. These techniques are attractive when data or computation is limited and may be interpretable through their transformations. However, a global adjustment can over-amplify noise, clip channels, or impose colors that are visually vivid but not faithful to the scene. Their success also varies with the degradation pattern because underwater effects are often spatially non-uniform.",
        "Physics-informed methods estimate quantities such as background light, transmission, attenuation, or depth and then invert an image-formation model. The dark channel prior inspired influential dehazing work [7], but underwater adaptations must address wavelength-dependent attenuation and cases where assumptions about dark pixels or illumination fail. Underwater-specific priors and color compensation methods attempt to account for those differences [6], [8]. Such approaches can provide a meaningful model of degradation, yet errors in estimated transmission or ambient light can create halos, color imbalance, and amplified noise. A useful method must therefore balance physical interpretation with robustness to scenes that violate its assumptions.",
        "Supervised deep networks learn a mapping from degraded inputs to reference targets. U-Net-like regressors, residual models, adversarial systems, and later generative approaches can capture nonlinear corrections from data. The UIEB work supplied a benchmark and introduced Water-Net, which uses gated fusion of transformed inputs [4], while GAN-based approaches such as FUnIE-GAN emphasized efficient enhancement and unpaired or weakly paired practical settings [10]. These examples show the breadth of learned underwater processing, but their objectives, datasets, and protocols differ. Direct numerical comparison with this project would be misleading without rerunning methods on the same fixed split.",
        "The present method uses diffusion rather than direct deterministic regression or adversarial training. Diffusion provides an explicit noising process and trains a denoiser across noise levels. It can model a restoration distribution and has proved adaptable to conditional inverse problems. Its costs include iterative sampling and sensitivity to the schedule and denoiser. In this project the purpose is not to establish diffusion as superior to all alternatives. The purpose is narrower: hold the diffusion formulation fixed and determine how an ordinary U-Net block design and a residual block design behave on the paired task.",
    ])
    subheading(doc, "2.3 U-Net Backbones")
    paragraphs(doc, [
        "U-Net was introduced for biomedical image segmentation, where precise localization must be combined with context [2]. Its encoder reduces spatial resolution while increasing representational width, and its decoder recovers resolution. Skip connections copy encoder features to corresponding decoder stages, allowing high-resolution information to bypass the bottleneck. This structure has become common in restoration and diffusion systems because a denoiser needs both local detail and a receptive field broad enough to interpret corruption. The skip paths are concatenative in the implementation used here, so decoder blocks receive both upsampled deep features and same-scale encoder features.",
        "The baseline ConditionalUNet is deliberately compact. It has two down blocks, one bottleneck block, and two up blocks. Each ConvBlock applies two 3 by 3 convolutions, GroupNorm, and SiLU activation. A projected time embedding is added after the first normalized convolution. Average pooling performs downsampling; bilinear interpolation performs upsampling; and a 1 by 1 convolution emits three noise channels. The noisy target and turbid condition are concatenated to form the six-channel input. This is a valid conditional denoiser but not an exhaustive modern U-Net with attention, many resolutions, or multiple blocks per level.",
        "Compactness provides practical advantages: low memory use, rapid debugging, and a clear baseline. It also creates a comparison issue when a proposed architecture adds blocks or projections. The baseline width of 32 yields approximately 527,363 parameters, substantially fewer than the residual model's approximately 886,371. A performance gain over that baseline can reflect the residual operation, the extra bottleneck, parameter count, or interactions among them. For this reason, the wider base-channel-42 U-Net is not an optional side result; it is required to interpret the evidence honestly.",
    ])
    subheading(doc, "2.4 Residual Learning")
    paragraphs(doc, [
        "Deep residual learning reformulates a block so that its transformed path estimates a residual relative to a shortcut [3]. When input and output dimensions match, the shortcut can be an identity; when channel dimensions differ, a projection aligns them. The block output combines the shortcut with the transformed signal before the final activation. This can improve gradient propagation and makes it easier for a block to preserve useful information while learning a correction. In restoration, where much of an input's spatial organization should survive, that inductive bias is intuitively relevant, although intuition alone does not establish an empirical advantage.",
        "The ResidualUNet used here keeps the broad encoder-decoder arrangement of the baseline. Its down, bottleneck, and up modules are ResidualBlock instances, and it includes two bottleneck residual blocks instead of one ordinary bottleneck block. Time embeddings enter the residual transformations. Average pooling, bilinear interpolation, and encoder-decoder concatenations are retained. Consequently, the experiment does not compare a plain ResNet classifier with U-Net. It compares two U-Net-shaped conditional denoisers whose internal processing differs, with the proposed version introducing residual block shortcuts and additional bottleneck processing.",
        "Residual learning should not be confused with predicting a restored-image residual directly. The diffusion network predicts Gaussian noise, and residual describes the internal block topology. Likewise, the outer U-Net skip connections and inner residual shortcuts serve different roles: the former move same-scale encoder features to decoder stages, while the latter add a block input to its transformed branch. Both can preserve information, but at different structural levels. The measured improvements in SSIM and Delta E are consistent with useful information preservation; they do not by themselves prove a specific causal mechanism inside the network.",
    ])
    subheading(doc, "2.5 Diffusion Models for Restoration")
    paragraphs(doc, [
        "Denoising diffusion probabilistic models define a forward Markov process that gradually adds Gaussian noise to data and a learned reverse process that removes it [1]. If alpha-bar denotes the cumulative product of retained-signal coefficients, a noisy sample at timestep t can be drawn directly from the clean target as the square root of alpha-bar times the target plus the square root of one minus alpha-bar times standard Gaussian noise. Training commonly samples a timestep uniformly, constructs that noisy target, and asks a network to predict the noise. This avoids simulating every preceding forward step for each update.",
        "Conditional restoration adds information about the observed degradation. In this implementation, the turbid image remains available as a condition at every denoising call. The noisy clear target and turbid condition are concatenated spatially, while the scalar timestep is represented by a sinusoidal embedding passed through a small multilayer perceptron. The model therefore knows both what conditioning observation it should follow and the severity of artificial diffusion noise it must remove. During reverse sampling, the trained noise estimates are used across the configured schedule to produce a restored output.",
        "Diffusion has been adapted to image-to-image translation and inverse problems in several ways. Palette demonstrated a unified conditional diffusion approach to image-to-image tasks [11], and SR3 showed iterative refinement for super-resolution [12]. These works motivate the general formulation but do not establish performance on this project's underwater dataset. The implementation here uses a 100-step linear beta schedule from 0.0001 to 0.02 and 100 sampling steps. It is intentionally simpler than systems using learned variance, attention-heavy denoisers, classifier-free guidance, accelerated samplers, or thousands of diffusion steps.",
        "Sampling cost is a practical limitation. A direct network can produce an output in one pass, whereas the implemented reverse process repeatedly invokes the denoiser. Reducing the number of steps may accelerate inference but can change image quality. Because the backbone comparison uses the same schedule and sampler, relative results are controlled with respect to this factor. They should not be interpreted as a complete exploration of the accuracy-speed frontier for underwater restoration.",
    ])
    subheading(doc, "2.6 Image Quality Metrics")
    paragraphs(doc, [
        "No single full-reference metric captures every desirable property. Mean squared error averages squared channel-wise pixel differences and strongly penalizes large deviations. Mean absolute error averages absolute differences and is less dominated by outliers. Peak signal-to-noise ratio is a logarithmic transformation related to MSE for a fixed data range; it is reported in decibels, with a higher value indicating lower pixel error. PSNR is convenient and widely recognizable, but equal MSE or PSNR can correspond to visually different distortions, and neither explicitly models structure or perception [13].",
        "Structural Similarity compares local luminance, contrast, and structure rather than treating all pixel errors independently [14]. SSIM is bounded in common use, and higher values indicate closer structural agreement. It is particularly relevant to the USAF-target motivation, where preservation of neighboring bars and boundaries matters. Nevertheless, SSIM remains a mathematical index with choices of window and implementation. A high value should not be equated automatically with perfect perceptual quality, physical correctness, or downstream-task utility.",
        "CIE76 Delta E measures Euclidean distance between colors represented in CIELAB space [15], [16]. Lower Delta E indicates a smaller average color difference under that representation. CIELAB was designed to be more perceptually organized than raw RGB, making Delta E useful for chart-like data where hue and lightness recovery matter. CIE76 is the simplest Delta E formula and is not perfectly perceptually uniform; later formulas such as CIEDE2000 account for known nonuniformities. The report therefore names the exact CIE76 variant rather than using the ambiguous label color error.",
        "Entropy summarizes the distribution of output intensity values and is included as a descriptive statistic. It is not a full-reference quality measure and has no universal better direction: high entropy may indicate retained detail or noise, while low entropy may indicate smoothness or lost texture. In this dataset, lower restored-output entropy accompanies the residual model's lower color error and higher SSIM, which is consistent with reduced noisy variation. That relationship is interpretive rather than definitive. The conclusions consequently rely most heavily on paired MSE, MAE, PSNR, SSIM, and Delta E, while treating entropy as supporting evidence.",
    ])


def add_system_requirements(doc: Document) -> None:
    heading(doc, "3 SYSTEM REQUIREMENTS")
    subheading(doc, "3.1 Design Constraints and Relevant Engineering Standards")
    paragraphs(doc, [
        "The first design constraint is data availability. Quantitative supervision requires aligned degraded and reference images, and only the 3,672 verified pairs support the measured restoration task. Auxiliary underwater scenes and video-derived frames lack corresponding clear targets. They may be inspected qualitatively but cannot enter full-reference metric tables. This separation is an evaluation constraint as well as a scientific-integrity requirement: assigning fabricated targets or comparing an external output with an unrelated image would produce meaningless numbers.",
        "Compute and storage impose a second constraint. Full experiments were conducted in Google Colab with A100 availability, while checkpoints, datasets, logs, and generated samples were kept outside the source repository. The selected 128-pixel main resolution, batch size 16, compact networks, and 100 diffusion steps balance feasibility with a meaningful architecture comparison. A 256-pixel ablation tests one higher-resolution setting but does not turn the system into an arbitrary-resolution production service. The implementation processes square RGB images and has no video-temporal component.",
        "Reproducibility requirements follow ordinary engineering and trustworthy-ML practice. Configurations record model names, widths, diffusion steps, image size, batch size, epochs, learning rate, seed, and artifact paths. Fixed filename lists define splits. The same test items are used for all table rows. Results are reported with configuration labels rather than selected anonymous runs. These practices are consistent with transparency and traceability principles emphasized in trustworthy AI guidance such as the NIST AI Risk Management Framework [17], although this student prototype is not a claim of formal framework compliance or safety certification.",
        "Image values are represented as RGB tensors normalized to [-1, 1], and evaluation outputs are returned to a bounded display range. Metric names include their direction and, for color, the exact Delta E variant. Dataset scope and limitations are disclosed. No claim is made that the model meets subsea operational safety, metrology, or autonomous-navigation standards. Such deployment would require camera calibration, environmental testing, robustness analysis, latency assessment, failure detection, and application-specific verification beyond the present study.",
    ])
    add_table_caption(doc, "Table 3-1: Principal design constraints and project responses")
    add_table(doc, [
        ["Constraint", "Project response", "Residual risk"],
        ["Paired data", "Use 3,672 verified filename pairs; fixed splits", "Limited scene diversity"],
        ["Unpaired external scenes", "Qualitative use only", "No cross-dataset full-reference score"],
        ["Compute", "Compact models; Colab GPU; 128-pixel main setting", "Iterative sampling remains costly"],
        ["Fair comparison", "Shared objective, split, schedule, and metrics", "Backbones still differ internally"],
        ["Capacity confound", "Add approximately parameter-matched U-Net", "Not a complete architecture sweep"],
        ["Deployment", "Research prototype only", "No real-time or safety certification"],
    ], font_size=9)
    subheading(doc, "3.2 Functional Requirements")
    paragraphs(doc, [
        "The data subsystem shall discover clear and turbid image directories, read a prescribed split file, and return corresponding items by the same filename. It shall reject a missing or empty split file and raise an error if a requested paired image is absent. Each item shall contain a normalized clear tensor, normalized turbid tensor, and filename. Optional horizontal augmentation shall apply the same flip to both members of a pair so spatial correspondence is never broken.",
        "The training subsystem shall sample a valid diffusion timestep for every item in a batch, generate a noisy target using the schedule, pass the noisy target, turbid condition, and timesteps to the selected model, and minimize noise-prediction loss. It shall support at least the conditional U-Net and residual model through the model factory and configurable base width. Checkpoint and log paths shall be independent of source code so full Colab artifacts can remain in external storage.",
        "The evaluation subsystem shall restore every item in the fixed test list and aggregate MSE, MAE, PSNR, SSIM, CIE76 Delta E, and entropy. It shall record the number of evaluated examples. Qualitative output shall preserve enough ordering information to compare turbid input, restored result, and clear reference for the same filename. Report generation shall consume the frozen CSV rather than silently invent or round-trip new metrics.",
        "The reporting subsystem shall identify which panels are dataset conditions and which are generated model outputs. It shall show the default comparison, capacity control, duration ablation, and resolution ablation. It shall retain the official institutional cover, provide readable tables and compact figures, and state the central result without universal superiority claims. The executable report generator and source configurations provide the reproducible path from frozen artifacts to the final document.",
    ])
    subheading(doc, "3.3 Non-Functional Requirements")
    paragraphs(doc, [
        "Reproducibility is the leading non-functional requirement. Seed 42, fixed split files, explicit YAML configurations, deterministic filenames, and a single metric summary make comparisons auditable. Complete bitwise repeatability on different GPU software stacks is not guaranteed, but a future runner can identify the intended setup and rerun it. Maintainability is supported by separating data, models, diffusion, sampling, training, evaluation, and scripts rather than embedding all logic in one notebook.",
        "Correctness is protected by lightweight automated tests for pairing, data loading, diffusion behavior, model forward passes, metrics, and a training step. Shape checks matter because both models concatenate a noisy RGB tensor and a conditional RGB tensor, then must return exactly three channels at the original spatial resolution. Input validation covers invalid diffusion step counts, timestep shapes, missing files, and empty splits. These checks do not replace end-to-end visual inspection, but they reduce the likelihood of reporting results from a structurally incorrect pipeline.",
        "Usability and clarity apply to generated artifacts. A report reader should be able to connect an experiment label to resolution, duration, and model family without decoding directory names. Tables should fit the page at readable type sizes. Model grids should align the same sample across Turbid, Baseline restored, Residual restored, and Clear columns. Captions should state whether higher or lower is preferred and whether an image is an external condition or model output.",
        "Efficiency is bounded rather than optimized. Training must fit available GPU memory at batch size 16 in the main setting, and report generation must complete locally without requiring the dataset or checkpoints. Inference uses all 100 configured reverse steps, favoring comparability over latency. Reliability outside the test distribution is explicitly limited; the program should not present an external qualitative result as guaranteed correction. These non-functional requirements define a transparent research prototype rather than a production underwater-vision product.",
    ])
    subheading(doc, "3.4 Evaluation Methodology")
    paragraphs(doc, [
        "The primary evaluation unit is the image pair in the 368-item test split. Training and validation examples are excluded from final metric aggregation. The frozen summary reports n=368 for every experiment, enabling direct row-wise comparison. The main contrast is default U-Net versus residual at 128 pixels and 50 epochs. Because all surrounding choices are shared, this contrast estimates the combined effect of adopting the implemented residual backbone design rather than changing the data or diffusion objective.",
        "A second contrast compares residual with the parameter-matched U-Net. Their parameter counts differ by about 15,426, or less than two percent relative to the residual count, while the default U-Net is much smaller. This control asks whether width alone can recover performance. It is not perfectly controlled for computation, activation count, receptive behavior, or block depth, and it does not isolate every residual component. It nevertheless provides substantially stronger evidence than comparing only the unequal default models.",
        "Two ablations probe training duration and input resolution. The 100-epoch rows are compared with their 50-epoch family counterparts at 128 pixels. The 256-pixel rows are compared with 128-pixel, 50-epoch rows. These are one-factor practical probes, not exhaustive hyperparameter optimizations. A metric that worsens after longer training or higher resolution may indicate optimization mismatch, stochastic variation, or a trade-off; it does not prove that the factor can never help.",
        "Interpretation uses direction-aware metrics and avoids selecting only favorable numbers. Lower MSE, MAE, and Delta E and higher PSNR and SSIM are preferred. Entropy is described, not treated as a universal objective. Qualitative grids use matching filenames and shared rows so visual comparison is legitimate. External source-condition figures are explicitly separated from model-output figures. Conclusions are restricted to the observed dataset, architecture implementations, schedules, and checkpoints.",
    ])


def add_system_architecture(doc: Document) -> None:
    heading(doc, "4 SYSTEM ARCHITECTURE")
    subheading(doc, "4.1 Dataset and Pairing Pipeline")
    paragraphs(doc, [
        "The architecture begins with two directories containing clear and turbid RGB files and text files containing split filenames. UnderwaterPairedDataset reads non-empty lines from one split file and uses each filename to open one image from each directory. This filename-driven design is preferable to independently sorting two directories because it makes the relationship explicit and fails loudly when one member is missing. The returned dictionary includes clear, turbid, and filename fields, preserving traceability through loading and visualization.",
        "Both images undergo the same deterministic resize to the configured square dimensions, tensor conversion, and channel-wise normalization with mean 0.5 and standard deviation 0.5. The resulting range is approximately [-1, 1], matching the model and diffusion operations. If augmentation is enabled, one random horizontal-flip decision is shared by both images. Applying independent geometric transforms would destroy pixel alignment and invalidate full-reference training, so paired augmentation is an architectural requirement rather than an incidental implementation detail.",
        "Fixed training, validation, and test lists sit outside the Dataset class. This separation allows the same loader code to serve each phase while keeping assignment stable across experiments. Training updates parameters, validation tracks behavior during development, and testing is reserved for final aggregate comparison. The counts sum exactly: 2,937 plus 367 plus 368 equals 3,672. No external unpaired item enters these split counts.",
    ])
    add_picture(doc, REPORT_ASSETS / "candidate_pair_sheet.png", 6.3)
    caption(doc, "Figure 4-1: Candidate turbid-clear pairs used to verify correspondence before training")
    subheading(doc, "4.2 Conditional Diffusion Pipeline")
    paragraphs(doc, [
        "DiffusionSchedule constructs 100 beta values linearly spaced from 0.0001 to 0.02. Alpha at each step is one minus beta, and the cumulative product of alphas determines how much clean signal remains. For a batch, diffusion_training_step samples one integer timestep per item, gathers the appropriate cumulative coefficients, and combines the clear tensor with newly sampled Gaussian noise. It returns both the noisy target and exact noise, so the supervised objective has a known target even though the reverse restoration target is generated iteratively.",
        "The model call has three inputs: noisy clear-domain tensor, turbid condition tensor, and timestep vector. Concatenation of the image tensors provides six spatial channels. A sinusoidal encoding maps each scalar timestep into a 128-dimensional representation; linear layers and SiLU transform it before block-specific projections inject it into feature maps. The output shape matches the three-channel noise tensor. Noise-prediction MSE is backpropagated, gradients update the model, and the scalar loss is recorded.",
    ])
    add_picture(doc, REPORT_ASSETS / "proposed_pipeline.png", 6.8)
    caption(doc, "Figure 4-2: Shared conditional diffusion pipeline with the denoising backbone as the experimental variable")
    paragraphs(doc, [
        "Figure 4-2 is a conceptual view. In the actual training path, artificial forward noise is applied to the clear target while the turbid input conditions noise prediction. At inference, the reverse process repeatedly evaluates the denoiser under the condition to construct a restored image. The baseline and proposed runs use the same schedule, target definition, loss, split, and metric implementation. Thus, the backbone is the planned independent variable in the principal experiment.",
        "The sampling architecture is computationally iterative. One restored batch requires repeated network calls rather than a single regression pass. That choice keeps the implementation close to the DDPM noise-prediction framework, but it has deployment implications. A real-time underwater vehicle might need fewer steps, a deterministic accelerated sampler, distillation, or a direct model. Those alternatives are beyond the controlled comparison and are left for future work.",
    ])
    subheading(doc, "4.3 Baseline Conditional U-Net")
    paragraphs(doc, [
        "The baseline starts with a ConvBlock mapping six channels to 32. After average pooling, a second block maps 32 to 64. Another pooling operation feeds a 128-channel bottleneck. The decoder upsamples bilinearly to the second encoder resolution, concatenates the 64-channel skip tensor, and processes the combined 192 channels into 64. It then upsamples to the first resolution, concatenates the 32-channel skip, and produces 32 channels before the final 1 by 1 output convolution.",
        "Each ConvBlock applies convolution, GroupNorm, SiLU, addition of a linearly projected time embedding, a second convolution, GroupNorm, and SiLU. Group counts are chosen from 8, 4, 2, or 1 according to channel divisibility. GroupNorm avoids dependence on batch statistics and is suitable for diffusion training where memory constraints can limit batch size. Average pooling and bilinear interpolation avoid learnable downsampling and transposed-convolution parameters, keeping the baseline compact and straightforward.",
        "The U-Net's long skip concatenations offer high-resolution evidence to the decoder, but the internal ConvBlocks do not add their inputs to outputs. The model has approximately 527,363 trainable parameters at base width 32. The parameter-matched variant changes only base width to 42 in the configuration and reaches approximately 901,797 parameters. Its role is to preserve the same ordinary block family while bringing capacity close to the residual model.",
    ])
    subheading(doc, "4.4 Proposed Residual Backbone")
    paragraphs(doc, [
        "The proposed ResidualUNet accepts the same six-channel input, uses the same time-embedding dimensionality, and follows the same 32, 64, and 128 channel scales. Its encoder and decoder modules are ResidualBlocks. It also uses two 128-channel bottleneck residual blocks, adding processing at the lowest spatial scale. A final 1 by 1 convolution maps 32 decoder channels to the three-channel noise estimate. The outer encoder-decoder skip concatenations, pooling, and interpolation remain aligned with the baseline design.",
        "Within each residual block, a transformed branch applies normalization, activation, convolution, and timestep conditioning while a shortcut provides an identity or channel-aligning projection. The combination allows information and gradients to bypass the nonlinear branch. When a block maps concatenated decoder input to fewer output channels, the projection is necessary. These projections and the additional bottleneck explain much of the parameter increase relative to the default U-Net.",
        "The proposed model has approximately 886,371 trainable parameters. It is slightly smaller than the capacity-control U-Net but substantially larger than the default baseline. Table 4-1 reports these counts openly. Parameter count is not equivalent to computational cost or expressive behavior, but it is an important first-order confound. The experiment therefore supports two conclusions at different levels: residual is clearly better than the project's default small U-Net on most reported quantities, and residual retains advantages in SSIM, Delta E, and entropy when compared at roughly equal parameter count.",
    ])
    add_table_caption(doc, "Table 4-1: Trainable parameter counts for the principal backbone configurations")
    add_table(doc, [
        ["Model", "Base channels", "Approx. parameters", "Role"],
        ["Default conditional U-Net", "32", "527,363", "Main baseline"],
        ["Residual U-Net", "32", "886,371", "Proposed backbone"],
        ["Parameter-matched U-Net", "42", "901,797", "Capacity control"],
    ], font_size=10)


def add_dataset_and_development(doc: Document) -> None:
    heading(doc, "5 DATASET PREPARATION AND MODEL DEVELOPMENT")
    subheading(doc, "5.1 Dataset Summary")
    paragraphs(doc, [
        "The supervised dataset consists of 3,672 matched observations of underwater color patches. Each pair contains a turbid image and a clearer reference image with the same filename. Color-chart data is useful for studying attenuation and cast because the patches provide repeated, spatially organized colors. The turbid images visibly lose contrast and depart from their references. This controlled structure also makes CIELAB color difference informative in addition to generic reconstruction metrics.",
        "Dataset preparation included checking directory contents, matching filenames, and visualizing candidate pairs. The main example sheet in Figure 5-1 places turbid and clear observations together so swapped labels, missing partners, gross misalignment, or anomalous files can be noticed before expensive training. Visual verification complements automated pairing: matching names guarantee a declared relationship but cannot prove that file contents are correct.",
    ])
    add_picture(doc, REPORT_ASSETS / "main_dataset_example_sheet.png", 6.4)
    caption(doc, "Figure 5-1: Representative paired training-domain examples; each turbid patch is associated with a clear reference")
    paragraphs(doc, [
        "The dataset is divided into 80 percent training and approximately 10 percent each for validation and testing: 2,937, 367, and 368 pairs. Integer counts reflect the finite collection. The same lists are reused across every experiment in the frozen results. This is more important than exact percentage labels because repeated random splitting could change sample difficulty and make small metric differences incomparable.",
        "External real scenes, the three-condition chart, and the USAF target serve different reporting purposes. Condition panels document visible underwater degradation. Auxiliary unpaired images or video-derived frames probe qualitative behavior outside the patch distribution. Neither category expands the paired training count, and neither receives invented full-reference metrics. Keeping those roles distinct prevents a visually interesting external panel from being mistaken for evidence of measured generalization.",
    ])
    add_table_caption(doc, "Table 5-1: Fixed paired-dataset split used by every quantitative experiment")
    add_table(doc, [
        ["Split", "Pairs", "Share", "Use"],
        ["Training", "2,937", "79.98%", "Parameter optimization"],
        ["Validation", "367", "9.99%", "Development monitoring"],
        ["Test", "368", "10.02%", "Frozen final evaluation"],
        ["Total", "3,672", "100.00%", "Paired supervised collection"],
    ], font_size=10)
    subheading(doc, "5.2 Preprocessing and Fixed Splits")
    paragraphs(doc, [
        "PIL opens each file and converts it to RGB, eliminating ambiguity from grayscale, palette, or alpha-channel inputs. Torchvision resizes both pair members to the configured dimensions with antialiasing. ToTensor maps display values to [0, 1], and normalization with mean and standard deviation 0.5 maps them to approximately [-1, 1]. The denormalization utility reverses this operation and clamps outputs to [0, 1] for visualization and metric handling as appropriate.",
        "Resizing standardizes batches but can alter fine detail. At 128 pixels, small structures may be smoothed or compressed; at 256 pixels, more spatial samples are available but optimization and memory demands increase. The paired chart patches are compatible with square resizing, yet the transformation would distort arbitrary non-square scenes unless aspect-ratio-aware preprocessing were added. The resolution ablation evaluates the implemented choice rather than claiming preservation of native camera geometry.",
        "Horizontal flipping is the available training augmentation. The same random decision applies to clear and turbid images, maintaining correspondence. No random color jitter is used because independent or aggressive color transformations could change the target relationship central to color restoration. No synthetic turbidity generation is included in the frozen setup. As a result, the model learns from the observed paired conditions rather than a parametrically expanded family of water effects.",
        "Split files are read once when a Dataset instance is created, and blank lines are ignored. A missing split raises FileNotFoundError; an empty one raises ValueError. During retrieval, either missing member also raises FileNotFoundError. These fail-fast behaviors protect experiments from silently shortening a split or pairing different arrays by index. The reported n=368 for all metric rows provides a final aggregate check that the complete test split was evaluated.",
    ])
    subheading(doc, "5.3 Training Configuration")
    paragraphs(doc, [
        "The main configurations use image size 128, batch size 16, 50 epochs, learning rate 0.0001, and random seed 42. The diffusion schedule contains 100 timesteps and sampling also uses 100 steps. Both model families output three channels from a six-channel concatenated input. The default and residual configurations differ in model name while retaining base width 32. The capacity-control file selects the ordinary conditional U-Net and raises base width to 42.",
        "At each update, a batch of clear and turbid tensors moves to the selected device. One timestep is sampled independently for each item. The schedule creates a noisy clear tensor and returns the exact Gaussian noise. The denoiser predicts that noise; mean squared noise-prediction loss is differentiated; optimizer gradients are cleared; backpropagation runs; and parameters update. The method trains neither an adversarial discriminator nor an explicit perceptual or color loss. Improvements in SSIM or Delta E therefore emerge under the shared noise objective rather than direct optimization of those metrics.",
        "Full runs took place in Google Colab with A100 availability. Dataset archives, checkpoints, logs, and evaluation outputs remained in Drive-oriented paths to avoid placing large binary artifacts in source control. The repository contains the reusable implementation, configurations, notebook workflow, tests, and summary artifacts. This separation is practical but means reproducing training requires access to the data and adequate compute; reproducing the report itself requires only the frozen CSV and figure assets.",
        "The 100-epoch configurations extend duration while preserving 128-pixel resolution and model family. The 256-pixel configurations preserve 50 epochs while changing resolution. These comparisons were selected to investigate common improvement hypotheses: train longer and provide more spatial detail. Their results show why such changes should be measured rather than assumed. Longer training benefits the baseline more consistently, and higher resolution changes structural and pixel metrics differently.",
    ])
    add_table_caption(doc, "Table 5-2: Shared main training and diffusion configuration")
    add_table(doc, [
        ["Item", "Value", "Applies to"],
        ["Main image size", "128 x 128", "Baseline, residual, capacity control"],
        ["Batch size", "16", "Main configuration"],
        ["Epochs", "50", "Main configuration"],
        ["Learning rate", "0.0001", "All listed YAML main configurations"],
        ["Seed", "42", "All listed YAML main configurations"],
        ["Diffusion timesteps", "100", "All listed YAML main configurations"],
        ["Sampling steps", "100", "All listed YAML main configurations"],
        ["Input/output", "6 concatenated / 3 predicted", "Both backbone families"],
    ], font_size=9)
    subheading(doc, "5.4 Model Development Timeline")
    paragraphs(doc, [
        "Development proceeded through checkpoints rather than beginning with all ablations at once. The first stage established the paired dataset, filename checks, fixed split logic, and visual examples. The second added the Dataset/DataLoader path and diffusion schedule sanity checks. These stages reduced the risk that later model behavior would be interpreted before data correspondence and noising behavior were verified.",
        "The third stage implemented and debugged the default conditional U-Net. A forward pass test confirmed output shape, and a training-step test checked that the diffusion objective could update the model. The fourth introduced ResidualBlock and ResidualUNet, then evaluated the main 128-pixel, 50-epoch comparison. This sequence kept the residual change localized rather than combining it with a new objective or data pipeline.",
        "The fifth stage consolidated full test metrics and added duration and resolution experiments. The sixth added the parameter-matched U-Net and reviewed the limits of claims on unpaired external data. Capacity control changed the final interpretation in an important way: the residual model no longer appeared categorically best once U-Net width increased, even though it retained structural and color-oriented advantages. Reporting evolved accordingly from a simple proposed-versus-baseline statement to a qualified architecture-and-capacity analysis.",
        "This iterative timeline is an engineering result in itself. Early automated checks protect basic correctness; fixed evaluation artifacts protect consistency; and later controls test alternative explanations. The final report uses frozen values from metric_summary.csv rather than manually transcribing selected notebook cells. Table 5-3 summarizes the progression without presenting unrecorded intermediate scores.",
    ])
    add_table_caption(doc, "Table 5-3: Experiment and implementation progression")
    add_table(doc, [
        ["Stage", "Implemented work", "Decision enabled"],
        ["1 Dataset setup", "Pair checks, split creation, previews", "Establish valid supervision"],
        ["2 Diffusion pipeline", "Loader, noising, schedule tests", "Verify conditional training inputs"],
        ["3 Baseline", "Conditional U-Net and debug training", "Create reference architecture"],
        ["4 Proposed model", "Residual blocks and main comparison", "Test residual backbone"],
        ["5 Full evaluation", "Metrics, 100-epoch and 256-pixel runs", "Study duration/resolution"],
        ["6 Controls", "Matched U-Net and external-data scope review", "Separate capacity and scope claims"],
    ], font_size=9)


def add_results_and_evaluation(doc: Document) -> None:
    heading(doc, "6 RESULTS AND EVALUATION")
    subheading(doc, "6.1 Main Quantitative Results")
    paragraphs(doc, [
        "Table 6-1 contains every frozen quantitative result. Each row represents the same 368-item test split. MSE, MAE, and CIE76 Delta E are better when lower; PSNR and SSIM are better when higher. Entropy is included descriptively. Reporting all seven configurations together prevents the main claim from hiding duration or resolution outcomes that are less favorable to the proposed model.",
    ])
    add_table_caption(doc, "Table 6-1: Full metrics on the fixed 368-pair test split")
    add_metric_summary_table(doc)
    paragraphs(doc, [
        "The default 128/50 U-Net has the highest MSE and MAE among the 128-pixel rows and the lowest 128-pixel SSIM. The main residual row improves every listed value relative to that default row when lower entropy is interpreted as smoother output. However, the matched U-Net records the lowest overall MSE and MAE and the highest overall PSNR in this set. The 256-pixel residual records the highest SSIM, while the 128/50 residual records the lowest Delta E and entropy. There is therefore no single configuration that dominates every criterion.",
        "Differences across metrics are expected because they reward different image properties. Pixel losses favor close channel values at exact coordinates. SSIM rewards local structural agreement. Delta E evaluates distance after conversion to a color representation designed around perceptual organization. A model can reduce local structural distortion without minimizing every RGB error, or slightly improve PSNR while leaving a larger color-space discrepancy. The correct reading is a profile of behavior, not a one-dimensional league table.",
    ])
    add_picture(doc, REPORT_OUTPUTS / "training_loss_curves.png", 6.2)
    caption(doc, "Figure 6-1: Recorded training and validation loss curves for completed experiments")
    paragraphs(doc, [
        "The loss curves provide optimization context but do not replace test metrics. Training minimizes noise-prediction MSE at random timesteps, whereas the final metrics compare sampled restored images with clear references. A lower training loss does not guarantee proportional improvement in SSIM or Delta E. Curves can reveal convergence, instability, or a plateau, but model selection and conclusions must remain tied to the fixed test protocol and generated outputs.",
    ])
    subheading(doc, "6.2 Default Baseline vs Residual")
    paragraphs(doc, [
        "The principal same-configuration comparison is favorable to the residual model. MSE falls from 0.039444 to 0.035727, a 9.42 percent reduction. MAE falls from 0.168026 to 0.159895, a 4.84 percent reduction. PSNR increases by 0.622892 dB, from 15.681501 to 16.304393. SSIM increases from 0.610758 to 0.788963, a relative increase of 29.18 percent. CIE76 Delta E falls from 29.645203 to 26.840257, a 9.46 percent reduction. Entropy falls from 5.095362 to 4.368050.",
        "SSIM is the largest relative change among the primary paired metrics. On these chart patches, structural consistency includes preservation of patch boundaries and within-patch organization. The lower Delta E indicates that the residual outputs are closer to the clear targets in CIELAB distance on average. Together, these outcomes support the claim that the residual backbone improves structural and color-oriented restoration behavior over the project's default conditional U-Net.",
        "The comparison does not identify residual shortcuts as the sole cause. ResidualUNet has more parameters and an extra bottleneck block, so it changes capacity and depth as well as shortcut topology. The result is best phrased as an improvement from replacing the default denoising backbone with the implemented residual backbone. The matched control in Section 6.3 then tests whether increasing ordinary U-Net width offers an alternative explanation.",
    ])
    add_picture(doc, REPORT_OUTPUTS / "ssim_comparison.png", 5.8)
    caption(doc, "Figure 6-2: SSIM across completed configurations; higher is better")
    add_picture(doc, REPORT_OUTPUTS / "delta_e_cie76_comparison.png", 5.8)
    caption(doc, "Figure 6-3: CIE76 Delta E across completed configurations; lower is better")
    add_picture(doc, REPORT_OUTPUTS / "compact_default_comparison.png", 6.4)
    caption(doc, "Figure 6-4: Matching test samples arranged as Turbid, default U-Net restoration, residual restoration, and Clear reference")
    paragraphs(doc, [
        "Figure 6-4 is generated programmatically from matching baseline and residual grids. Only three representative rows are retained so every column remains legible and the figure fits within one page. Turbid and Clear come from the baseline sheet, while the two restored columns come from their corresponding model sheets. The arrangement supports direct inspection of color drift, smoothness, and patch boundaries without asking the reader to compare separate tall figures.",
    ])
    subheading(doc, "6.3 Parameter-Matched Capacity Control")
    paragraphs(doc, [
        "The capacity control is the most important qualification. Widening the ordinary U-Net from 32 to 42 base channels raises its parameter count from approximately 527,363 to 901,797, close to the residual model's approximately 886,371. At the same 128-pixel, 50-epoch setting, the matched U-Net obtains MSE 0.034294, MAE 0.155805, and PSNR 16.356273. These are slightly better than residual values of 0.035727, 0.159895, and 16.304393. The PSNR margin is only about 0.052 dB, but its direction is clear.",
        "The residual model remains stronger on SSIM, 0.788963 versus 0.746291, and Delta E, 26.840257 versus 27.225774. It also has lower entropy, 4.368050 versus 4.790907. This pattern suggests that width and residual topology distribute performance differently across criteria. Increased U-Net capacity is effective for direct reconstruction error, while the residual configuration provides the best structural and color-space scores at this resolution and duration.",
        "A capacity-matched result is not a proof of architecture in isolation. Parameter count does not match depth, operation count, memory traffic, or optimization geometry, and base width 42 is only one possible U-Net design. Multiple seeds would be needed to estimate run-to-run variance. Still, the control falsifies an overly broad narrative that the residual model is best on every metric and provides evidence that the default baseline was partly capacity limited.",
    ])
    add_table_caption(doc, "Table 6-2: Focused 128-pixel, 50-epoch capacity-control comparison")
    add_table(doc, [
        ["Metric", "Matched U-Net", "Residual", "Observed leader"],
        ["Parameters", "901,797", "886,371", "Approximately matched"],
        ["MSE", "0.034294", "0.035727", "Matched U-Net"],
        ["MAE", "0.155805", "0.159895", "Matched U-Net"],
        ["PSNR (dB)", "16.356273", "16.304393", "Matched U-Net"],
        ["SSIM", "0.746291", "0.788963", "Residual"],
        ["CIE76 Delta E", "27.225774", "26.840257", "Residual"],
        ["Entropy", "4.790907", "4.368050", "Residual (lower here)"],
    ], font_size=9)
    add_picture(doc, REPORT_OUTPUTS / "compact_capacity_comparison.png", 6.4)
    caption(doc, "Figure 6-5: Matching test samples for the parameter-matched U-Net and residual capacity comparison")
    subheading(doc, "6.4 Training Duration Ablation")
    paragraphs(doc, [
        "Extending the default U-Net from 50 to 100 epochs improves all of its listed metrics: MSE decreases from 0.039444 to 0.037111, MAE from 0.168026 to 0.163163, PSNR rises from 15.681501 to 15.976561, SSIM rises from 0.610758 to 0.623376, Delta E falls from 29.645203 to 28.571008, and entropy falls from 5.095362 to 4.943213. The baseline therefore had additional optimization benefit available under the unchanged longer schedule.",
        "The residual model does not improve consistently with another 50 epochs. MSE changes from 0.035727 to 0.035936 and MAE from 0.159895 to 0.160301, both slightly worse. PSNR rises from 16.304393 to 16.352538, while SSIM falls from 0.788963 to 0.785447, Delta E rises slightly from 26.840257 to 26.866293, and entropy rises from 4.368050 to 4.710935. These small mixed movements suggest a plateau or objective-specific trade-off rather than a clear gain.",
        "The appropriate inference is not that residual networks should never train longer. The learning rate, decay policy, regularization, checkpoint selection, and stochastic variation may influence the 100-epoch outcome. The frozen experiment shows only that doubling epochs under the current setup did not improve the residual model across the selected metrics, whereas it helped the default baseline. Future longer runs should tune the schedule rather than merely repeat identical epochs.",
    ])
    add_picture(doc, REPORT_OUTPUTS / "compact_duration_comparison.png", 6.4)
    caption(doc, "Figure 6-6: Matching 100-epoch U-Net and residual restorations with turbid inputs and clear references")
    subheading(doc, "6.5 Resolution Ablation")
    paragraphs(doc, [
        "At 256 pixels and 50 epochs, the baseline U-Net records MSE 0.039126, MAE 0.167594, PSNR 15.749889, SSIM 0.609380, Delta E 29.117468, and entropy 4.873310. Relative to its 128-pixel row, pixel errors and PSNR improve slightly, SSIM decreases slightly, Delta E improves, and entropy decreases. The changes are modest, so simply doubling width and height does not transform baseline behavior.",
        "The 256-pixel residual row reaches SSIM 0.807967, the highest SSIM in the summary. It also records PSNR 15.900427, MSE 0.039173, MAE 0.168835, Delta E 27.078848, and entropy 4.667301. Compared with residual 128/50, SSIM improves but MSE, MAE, PSNR, Delta E, and entropy worsen. Higher resolution therefore favors the structural index while failing to improve the broader metric profile under this training budget.",
        "Increasing resolution quadruples pixel count and changes optimization demands. A fixed number of epochs does not imply equal compute per pixel or equal effective convergence. The architecture has the same number of parameters but feature maps are larger, and fine-scale variation can affect error. A stronger resolution study would tune batch size, learning schedule, model receptive field, and sampling cost. The present row is correctly described as an ablation, not an optimized high-resolution model.",
    ])
    add_picture(doc, REPORT_OUTPUTS / "compact_resolution_comparison.png", 6.4)
    caption(doc, "Figure 6-7: Three matching 256-pixel examples comparing baseline and residual restorations")
    add_picture(doc, REPORT_OUTPUTS / "mse_comparison.png", 5.8)
    caption(doc, "Figure 6-8: MSE across completed configurations; lower is better")
    add_picture(doc, REPORT_OUTPUTS / "psnr_comparison.png", 5.8)
    caption(doc, "Figure 6-9: PSNR across completed configurations; higher is better")
    subheading(doc, "6.6 Qualitative and Cross-Dataset Scope")
    paragraphs(doc, [
        "Qualitative analysis complements aggregate scores by exposing the form of errors. In matching patch grids, useful questions include whether the restored patch approaches reference hue, whether boundaries remain coherent, whether noise or mottling is introduced, and whether the method collapses different turbid inputs toward a common output. The compact four-column figures make these questions easier to inspect than separate vertical model sheets. They also preserve filename-level correspondence.",
        "The paired test examples remain in distribution: they come from the same collection and task definition as training, although their filenames were held out. They are the proper visual counterpart to the numerical tables. A few displayed rows cannot represent all 368 examples and should not be selected as proof of average performance. Their role is illustrative, while the aggregate metrics carry the test-set claim.",
        "External natural underwater scenes and video-derived frames are more challenging because the training set is dominated by controlled color patches. Such images may contain objects, long-range haze, non-uniform depth, suspended particles, natural textures, and camera motions absent from the training distribution. The model can be run on them, but no paired clear reference is available. Consequently, visual plausibility is the only evidence available, and an apparently pleasing output could still alter true colors or remove meaningful detail.",
        "The three-condition chart and USAF images shown in the Introduction are not model outputs and should not be used to claim successful restoration. They demonstrate real degradation modes and motivate color and structure measurements. This distinction is especially important in generative restoration because a model can produce a plausible-looking image that is not faithful to the observed scene. A future external benchmark should use genuinely paired or carefully calibrated references before numerical generalization claims are made.",
        "Across the available evidence, the residual outputs tend to be structurally consistent and less noisy in the entropy sense, but they do not fully recover all clear-reference colors. The modest absolute PSNR and Delta E values reinforce that the problem remains unsolved. The project establishes a relative improvement under a specific protocol, not perfect de-turbiding or guaranteed recovery of information lost in the water medium.",
    ])
    subheading(doc, "6.7 Limitations")
    paragraphs(doc, [
        "The first limitation is dataset scope. Controlled color-patch pairs make alignment and color measurement possible, but they underrepresent the spatial complexity of natural underwater scenes. A model can exploit repeated chart layout and acquisition statistics. Performance on the fixed test split therefore demonstrates in-distribution paired restoration and should not be generalized automatically to coral reefs, pipelines, organisms, archaeological sites, or open-water robotics.",
        "The second limitation is experimental breadth. Results come from one frozen split and reported configurations, without multiple random seeds or confidence intervals. Small differences, such as the 0.052 dB PSNR gap between matched U-Net and residual, may be sensitive to training variability. The capacity control matches parameter count approximately but not computation, depth, or every architectural feature. The residual implementation changes block topology and bottleneck depth together, so component-level causal attribution is limited.",
        "The third limitation concerns objectives and metrics. Training uses only noise-prediction MSE; no explicit physical, perceptual, structural, or color-consistency term is optimized. Full-reference metrics reduce complex visual behavior to averages and may not align perfectly with human judgment or downstream tasks. CIE76 is a basic color-difference formula; entropy has no universal quality direction; and none of the metrics verifies physical scene radiance. Human evaluation and task-based evaluation are absent.",
        "The fourth limitation is practical deployment. Reverse diffusion uses 100 network evaluations, and runtime or energy measurements were not part of the frozen results. The model works on single square images and does not enforce temporal consistency, so independent video frames may flicker. There is no uncertainty estimate, out-of-distribution detector, or safeguard against hallucinated detail. Inputs from different cameras, water types, or lighting conditions may fall outside the learned mapping.",
        "Finally, higher resolution and longer duration were tested only with limited schedules. The 256-pixel residual improves SSIM while degrading several other metrics; the 100-epoch residual plateaus. These observations may reflect insufficient tuning rather than fundamental limitations. Honest reporting requires retaining these mixed outcomes and treating further optimization as future work rather than retroactively selecting only the best metric from each run.",
    ])


def add_conclusions(doc: Document) -> None:
    heading(doc, "7 CONCLUSIONS AND FUTURE WORKS")
    paragraphs(doc, [
        "This project implemented a supervised conditional diffusion pipeline for underwater image restoration and used it to compare denoising backbones. The input condition is a turbid underwater color-patch image; the reference is its filename-matched clearer image. The project verified 3,672 pairs, fixed training, validation, and test assignments, implemented shared preprocessing and a 100-step linear diffusion schedule, and evaluated all completed configurations on the same 368-pair test set.",
        "The baseline conditional U-Net and proposed ResidualUNet share the same six-channel conditional input, timestep embedding concept, encoder-decoder scales, outer skip connections, noise-prediction objective, optimizer configuration, and principal training budget. The proposed design replaces ordinary ConvBlocks with ResidualBlocks and adds bottleneck processing. At 128 pixels and 50 epochs, this replacement improves all reported values relative to the default base-width-32 U-Net, including a 29.18 percent relative SSIM increase and a 9.46 percent Delta E reduction.",
        "The parameter-matched U-Net is essential to the final conclusion. Widening the ordinary U-Net to approximately 901,797 parameters makes it slightly better than the approximately 886,371-parameter residual model on MSE, MAE, and PSNR. The residual remains better on SSIM, CIE76 Delta E, and entropy. Thus, the evidence does not support saying that residual processing is universally superior. It supports saying that the implemented residual backbone improves structural similarity and color-difference behavior, while ordinary U-Net capacity is highly competitive for exact pixel reconstruction.",
        "The ablations reinforce the need for metric-specific interpretation. Doubling training duration improves the default baseline but does not consistently improve residual performance. Doubling resolution gives the residual model the highest SSIM but worsens several pixel and color measures relative to its 128-pixel run. Neither more epochs nor more pixels is an automatic solution without learning-rate, regularization, architecture, and sampling adjustments.",
        "The most immediate future work is repeated training with several seeds and statistical summaries. Confidence intervals or paired per-image tests would show whether small matched-capacity differences are stable. A component ablation should compare equal-depth ordinary and residual blocks, remove or retain the second bottleneck systematically, and measure multiply-accumulate operations, memory, sampling time, and parameter count. This would isolate shortcut effects more cleanly than parameter matching alone.",
        "Data expansion is equally important. Training should include diverse paired natural scenes spanning water types, depths, distances, artificial illumination, cameras, and object classes. UIEB or another recognized benchmark can provide external context when its protocol is followed, but paired and unpaired subsets must be treated correctly. Synthetic degradation may expand coverage if its physical assumptions are validated against real data. Camera color calibration and raw-image access could reduce unknown processing effects.",
        "Model work can investigate attention-enhanced U-Nets, learned or cosine schedules, fewer-step samplers, deterministic DDIM-like sampling, diffusion distillation, and direct restoration baselines. Comparisons against a deterministic U-Net regression model, a residual regression model, representative physics-based methods, and established underwater enhancement networks would determine whether diffusion's iterative cost is justified. Such comparisons must use the same data split and report runtime alongside image quality.",
        "Objective design can incorporate color-aware, structural, or perceptual terms, but each addition should be evaluated for faithfulness rather than visual vividness alone. CIEDE2000 could complement CIE76. Human preference studies, underwater-domain task performance, and calibrated chart measurements would broaden evaluation. For video, temporal conditioning or consistency losses are needed to avoid frame-to-frame flicker. An uncertainty or out-of-distribution mechanism should warn when an input differs substantially from training data.",
        "In conclusion, residual learning is a useful architectural direction for this conditional diffusion restoration task, particularly for structural and color-oriented criteria. The controlled capacity result makes that contribution more credible by defining where the advantage does and does not hold. The resulting system, artifacts, and report form a reproducible foundation for broader data, stronger controls, faster sampling, and deployment-oriented evaluation, while remaining explicit about the limited domain and the absence of universal restoration guarantees.",
    ])


def add_references(doc: Document) -> None:
    heading(doc, "8 REFERENCES")
    refs = [
        "[1] J. Ho, A. Jain, and P. Abbeel, 'Denoising Diffusion Probabilistic Models,' Advances in Neural Information Processing Systems, vol. 33, pp. 6840-6851, 2020.",
        "[2] O. Ronneberger, P. Fischer, and T. Brox, 'U-Net: Convolutional Networks for Biomedical Image Segmentation,' Medical Image Computing and Computer-Assisted Intervention, pp. 234-241, 2015. doi: 10.1007/978-3-319-24574-4_28.",
        "[3] K. He, X. Zhang, S. Ren, and J. Sun, 'Deep Residual Learning for Image Recognition,' Proceedings of the IEEE Conference on Computer Vision and Pattern Recognition, pp. 770-778, 2016. doi: 10.1109/CVPR.2016.90.",
        "[4] C. Li, C. Guo, W. Ren, R. Cong, J. Hou, S. Kwong, and D. Tao, 'An Underwater Image Enhancement Benchmark Dataset and Beyond,' IEEE Transactions on Image Processing, vol. 29, pp. 4376-4389, 2020. doi: 10.1109/TIP.2019.2955241.",
        "[5] J. S. Jaffe, 'Computer Modeling and the Design of Optimal Underwater Imaging Systems,' IEEE Journal of Oceanic Engineering, vol. 15, no. 2, pp. 101-111, 1990. doi: 10.1109/48.50695.",
        "[6] D. Akkaynak and T. Treibitz, 'A Revised Underwater Image Formation Model,' Proceedings of the IEEE Conference on Computer Vision and Pattern Recognition, pp. 6723-6732, 2018.",
        "[7] K. He, J. Sun, and X. Tang, 'Single Image Haze Removal Using Dark Channel Prior,' IEEE Transactions on Pattern Analysis and Machine Intelligence, vol. 33, no. 12, pp. 2341-2353, 2011. doi: 10.1109/TPAMI.2010.168.",
        "[8] P. Drews, E. Nascimento, F. Moraes, S. Botelho, and M. Campos, 'Transmission Estimation in Underwater Single Images,' Proceedings of the IEEE International Conference on Computer Vision Workshops, pp. 825-830, 2013.",
        "[9] C. Li, J. Guo, and C. Guo, 'Emerging From Water: Underwater Image Color Correction Based on Weakly Supervised Color Transfer,' IEEE Signal Processing Letters, vol. 25, no. 3, pp. 323-327, 2018.",
        "[10] M. J. Islam, Y. Xia, and J. Sattar, 'Fast Underwater Image Enhancement for Improved Visual Perception,' IEEE Robotics and Automation Letters, vol. 5, no. 2, pp. 3227-3234, 2020.",
        "[11] C. Saharia, W. Chan, H. Chang, C. A. Lee, J. Ho, T. Salimans, D. J. Fleet, and M. Norouzi, 'Palette: Image-to-Image Diffusion Models,' ACM SIGGRAPH, 2022. doi: 10.1145/3528233.3530757.",
        "[12] C. Saharia, J. Ho, W. Chan, T. Salimans, D. J. Fleet, and M. Norouzi, 'Image Super-Resolution via Iterative Refinement,' IEEE Transactions on Pattern Analysis and Machine Intelligence, vol. 45, no. 4, pp. 4713-4726, 2023.",
        "[13] Q. Huynh-Thu and M. Ghanbari, 'Scope of Validity of PSNR in Image/Video Quality Assessment,' Electronics Letters, vol. 44, no. 13, pp. 800-801, 2008. doi: 10.1049/el:20080522.",
        "[14] Z. Wang, A. C. Bovik, H. R. Sheikh, and E. P. Simoncelli, 'Image Quality Assessment: From Error Visibility to Structural Similarity,' IEEE Transactions on Image Processing, vol. 13, no. 4, pp. 600-612, 2004. doi: 10.1109/TIP.2003.819861.",
        "[15] Commission Internationale de l'Eclairage, Colorimetry, 3rd ed., CIE Publication 15:2004, Vienna, 2004.",
        "[16] G. Sharma, W. Wu, and E. N. Dalal, 'The CIEDE2000 Color-Difference Formula: Implementation Notes, Supplementary Test Data, and Mathematical Observations,' Color Research and Application, vol. 30, no. 1, pp. 21-30, 2005. doi: 10.1002/col.20070.",
        "[17] National Institute of Standards and Technology, Artificial Intelligence Risk Management Framework (AI RMF 1.0), NIST AI 100-1, 2023. doi: 10.6028/NIST.AI.100-1.",
        "[18] F. Iqbal and B. U. Toreyin, 'Underwater Turbid Image Restoration Using Diffusion Models,' project reference manuscript supplied for this project.",
    ]
    for ref in refs:
        reference(doc, ref)


def create_compact_comparison_assets() -> None:
    comparisons = [
        (
            COLAB_RESULTS / "baseline_full" / "baseline_comparison_grid.png",
            COLAB_RESULTS / "residual_full" / "residual_comparison_grid.png",
            REPORT_OUTPUTS / "compact_default_comparison.png",
        ),
        (
            COLAB_RESULTS / "param_matched_unet" / "results" / "baseline_full" / "baseline_comparison_grid.png",
            COLAB_RESULTS / "residual_full" / "residual_comparison_grid.png",
            REPORT_OUTPUTS / "compact_capacity_comparison.png",
        ),
        (
            COLAB_RESULTS / "baseline_100_full" / "baseline_comparison_grid.png",
            COLAB_RESULTS / "residual_100_full" / "residual_comparison_grid.png",
            REPORT_OUTPUTS / "compact_duration_comparison.png",
        ),
        (
            COLAB_RESULTS / "img256" / "results" / "baseline_full" / "baseline_comparison_grid.png",
            COLAB_RESULTS / "img256" / "results" / "residual_full" / "residual_comparison_grid.png",
            REPORT_OUTPUTS / "compact_resolution_comparison.png",
        ),
    ]
    for baseline_path, residual_path, output_path in comparisons:
        if baseline_path.exists() and residual_path.exists():
            combine_comparison_grids(baseline_path, residual_path, output_path)


def combine_comparison_grids(baseline_path: Path, residual_path: Path, output_path: Path) -> None:
    baseline = Image.open(baseline_path).convert("RGB")
    residual = Image.open(residual_path).convert("RGB")
    cell_width = baseline.width // 3
    row_height = baseline.height // 4
    label_height = max(20, row_height - cell_width)
    sample_height = row_height - label_height
    selected_rows = (0, 1, 2)
    header_height = 34
    canvas = Image.new("RGB", (cell_width * 4, header_height + sample_height * 3), "white")
    draw = ImageDraw.Draw(canvas)
    font = ImageFont.load_default()
    headers = ("Turbid", "Baseline restored", "Residual restored", "Clear")
    for column, text in enumerate(headers):
        left = column * cell_width
        box = draw.textbbox((0, 0), text, font=font)
        draw.text((left + (cell_width - (box[2] - box[0])) / 2, 11), text, fill="black", font=font)
    for output_row, source_row in enumerate(selected_rows):
        top = source_row * row_height + label_height
        bottom = top + sample_height
        tiles = (
            baseline.crop((0, top, cell_width, bottom)),
            baseline.crop((cell_width, top, cell_width * 2, bottom)),
            residual.crop((cell_width, top, cell_width * 2, bottom)),
            baseline.crop((cell_width * 2, top, cell_width * 3, bottom)),
        )
        for column, tile in enumerate(tiles):
            if tile.size != (cell_width, sample_height):
                tile = tile.resize((cell_width, sample_height), Image.Resampling.LANCZOS)
            canvas.paste(tile, (column * cell_width, header_height + output_row * sample_height))
    canvas.save(output_path)


def set_paragraph_text(paragraph, text: str, keep_style: bool = True) -> None:
    style = paragraph.style if keep_style else None
    alignment = paragraph.alignment
    paragraph.clear()
    if style is not None:
        paragraph.style = style
    paragraph.alignment = alignment
    if text:
        paragraph.add_run(text)


def paragraphs(doc: Document, items: list[str]) -> None:
    for text in items:
        para(doc, text)


def heading(doc: Document, text: str) -> None:
    if doc.paragraphs and doc.paragraphs[-1].text.strip():
        doc.add_page_break()
    text = re.sub(r"^\d+\s+", "", text)
    p = doc.add_heading(text, level=1)
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after = Pt(8)
    p.paragraph_format.keep_with_next = True


def preamble_heading(doc: Document, text: str) -> None:
    p = doc.add_paragraph()
    try:
        p.style = "Preamble Title"
    except KeyError:
        p.style = "Normal"
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    p.paragraph_format.left_indent = Inches(0.5)
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after = Pt(6)
    run = p.add_run(text)
    run.font.name = "Calibri"
    run.font.size = Pt(12)
    run.bold = False


def subheading(doc: Document, text: str) -> None:
    text = re.sub(r"^\d+\.\d+\s+", "", text)
    p = doc.add_heading(text, level=2)
    p.paragraph_format.space_before = Pt(10)
    p.paragraph_format.space_after = Pt(5)
    p.paragraph_format.keep_with_next = True


def para(doc: Document, text: str) -> None:
    p = doc.add_paragraph(style="Normal")
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    p.paragraph_format.left_indent = Inches(0.48)
    p.paragraph_format.right_indent = Inches(0.48)
    p.paragraph_format.first_line_indent = Inches(0)
    p.paragraph_format.space_after = Pt(5)
    p.paragraph_format.line_spacing = 1.08
    run = p.add_run(text)
    run.font.name = "Calibri"
    run.font.size = Pt(12)


def reference(doc: Document, text: str) -> None:
    p = doc.add_paragraph(style="Normal")
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    p.paragraph_format.left_indent = Inches(0.28)
    p.paragraph_format.first_line_indent = Inches(-0.28)
    p.paragraph_format.space_after = Pt(6)
    p.add_run(text)


def add_picture(doc: Document, path: Path, width: float) -> None:
    if not path.exists():
        para(doc, f"Missing figure asset: {path.name}")
        return
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(5)
    p.paragraph_format.space_after = Pt(2)
    p.paragraph_format.keep_with_next = True
    p.add_run().add_picture(str(path), width=Inches(width))


def caption(doc: Document, text: str) -> None:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(8)
    p.paragraph_format.keep_together = True
    run = p.add_run(text)
    run.font.name = "Calibri"
    run.font.size = Pt(10)
    run.italic = False


def add_table_caption(doc: Document, text: str) -> None:
    caption(doc, text)


def add_table(doc: Document, rows: list[list[str]], font_size: int = 9) -> None:
    table = doc.add_table(rows=len(rows), cols=max(len(row) for row in rows))
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.autofit = True
    table.rows[0]._tr.get_or_add_trPr().append(OxmlElement("w:tblHeader"))
    for row_index, row in enumerate(rows):
        for column_index, value in enumerate(row):
            cell = table.cell(row_index, column_index)
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            cell.text = value
            if row_index == 0:
                shading = OxmlElement("w:shd")
                shading.set(qn("w:fill"), "D9EAF7")
                cell._tc.get_or_add_tcPr().append(shading)
            for paragraph in cell.paragraphs:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER if row_index == 0 else WD_ALIGN_PARAGRAPH.LEFT
                paragraph.paragraph_format.space_after = Pt(0)
                for run in paragraph.runs:
                    run.font.name = "Calibri"
                    run.font.size = Pt(font_size)
                    run.bold = row_index == 0
    spacer = doc.add_paragraph()
    spacer.paragraph_format.space_after = Pt(0)


def add_metric_summary_table(doc: Document) -> None:
    rows = list(csv.DictReader((REPORT_OUTPUTS / "metric_summary.csv").open(encoding="utf-8")))
    names = {
        "baseline_128_50": "U-Net 128/50",
        "param_matched_unet_128_50": "Matched U-Net 128/50",
        "residual_128_50": "Residual 128/50",
        "baseline_128_100": "U-Net 128/100",
        "residual_128_100": "Residual 128/100",
        "baseline_256_50": "U-Net 256/50",
        "residual_256_50": "Residual 256/50",
    }
    table_rows = [["Experiment", "MSE", "MAE", "PSNR", "SSIM", "Delta E", "Entropy"]]
    for row in rows:
        table_rows.append([
            names.get(row["experiment"], row["experiment"]),
            row["mse"], row["mae"], row["psnr"], row["ssim"],
            row["delta_e_cie76"], row["entropy"],
        ])
    add_table(doc, table_rows, font_size=8)


def create_presentation_outline() -> None:
    text = """# Final Presentation Outline

## Slide 1 - Title
Enhancing Diffusion-Based Underwater Image Restoration with Residual Networks.

## Slide 2 - Problem
Underwater turbidity causes scattering, low contrast, and color distortion. The task is paired restoration from turbid image to clear/reference image.

## Slide 3 - Dataset
3672 paired images: 2937 train, 367 validation, 368 test. External/unpaired datasets are qualitative only.

## Slide 4 - Method
Conditional diffusion model with the same noising objective. Baseline uses a U-Net denoiser; proposed model uses a ResNet-style residual denoiser.

## Slide 5 - Main Result
Residual improves strongly over the default U-Net baseline, especially on SSIM and Delta E.

## Slide 6 - Capacity Control
Parameter-matched U-Net improves MSE/PSNR. Residual remains better on SSIM, Delta E, and entropy.

## Slide 7 - Qualitative Examples
Show residual comparison grid and parameter-matched U-Net grid.

## Slide 8 - Conclusion
Residual diffusion improves structural and color-oriented restoration, but the result is not a universal win on every metric.
"""
    (DELIVERABLES / "final_presentation_outline.md").write_text(text, encoding="utf-8")


def create_presentation_pptx() -> None:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    add_slide(prs, "Enhancing Diffusion-Based Underwater Image Restoration", ["Residual denoising backbones inside conditional diffusion", "Student ID: 150210321"])
    add_slide(prs, "Problem", ["Turbidity causes scattering, low contrast, and color distortion.", "Task: turbid underwater patch -> clear/reference patch.", "Question: does residual denoising improve the diffusion backbone?"])
    add_picture_slide(prs, "Dataset Example", REPORT_ASSETS / "paired_underwater_patches.png")
    add_slide(prs, "Method", ["Conditional diffusion model predicts noise added to the clean target.", "Baseline: conditional U-Net denoising backbone.", "Proposed: ResNet-style residual denoising backbone.", "Capacity control: parameter-matched U-Net."])
    add_picture_slide(prs, "Training Curves", REPORT_OUTPUTS / "training_loss_curves.png")
    add_slide(prs, "Main Result", ["Residual 128/50 improves default U-Net 128/50.", "SSIM: 0.610758 -> 0.788963.", "Delta E: 29.645203 -> 26.840257.", "MSE reduction: 9.42%."])
    add_picture_slide(prs, "SSIM Comparison", REPORT_OUTPUTS / "ssim_comparison.png")
    add_slide(prs, "Capacity Control", ["Matched U-Net is best on MSE, MAE, and PSNR.", "Residual is best on SSIM, Delta E, and entropy.", "This makes the final claim more honest and stronger."])
    add_picture_slide(prs, "Residual Examples", COLAB_RESULTS / "residual_full" / "residual_comparison_grid.png")
    add_slide(prs, "Conclusion", ["Residual backbone improves structural/color-oriented restoration over the default baseline.", "Parameter matching shows model capacity also matters.", "External datasets are qualitative only because paired references are unavailable."])
    prs.save(DELIVERABLES / "final_presentation_150210321.pptx")


def add_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for index, bullet_text in enumerate(bullets):
        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
        paragraph.text = bullet_text
        paragraph.level = 0
        paragraph.font.size = PptPt(24)


def add_picture_slide(prs: Presentation, title: str, image_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), PptInches(1.05), PptInches(1.15), width=PptInches(11.2))
    else:
        box = slide.shapes.add_textbox(PptInches(1.0), PptInches(2.5), PptInches(11.0), PptInches(1.0))
        box.text_frame.text = f"Missing image: {image_path}"


def create_submission_readme() -> None:
    text = """# Final Submission Notes

This folder contains the final local deliverables.

## Files

- `final_report_150210321.docx`: editable final report.
- `final_presentation_150210321.pptx`: editable final presentation.
- `final_presentation_outline.md`: short slide outline.
- `README_FINAL_SUBMISSION.md`: this file.

## GitHub

Repository: https://github.com/bedirhanozturk1/Underwater_Resnet_Restoration

## Reporting Note

Use the capacity-control interpretation. The residual backbone is strongest on SSIM, Delta E, and entropy. The parameter-matched U-Net is slightly stronger on MSE, MAE, and PSNR.
"""
    (DELIVERABLES / "README_FINAL_SUBMISSION.md").write_text(text, encoding="utf-8")


if __name__ == "__main__":
    main()
